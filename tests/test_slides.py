"""Tests for PPTX slide rendering and image-heavy detection."""

import tempfile
from pathlib import Path
from unittest.mock import patch, MagicMock

import pytest

from src.slides.renderer import (
    RenderedSlide,
    MAX_SLIDES,
    detect_image_heavy,
    can_render,
)


# ── detect_image_heavy tests ─────────────────────────────────────


class TestDetectImageHeavy:
    """Test the image-heavy detection heuristic."""

    def _make_pptx(self, num_slides: int, text_per_slide: str = "") -> Path:
        """Create a minimal PPTX file for testing."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        for i in range(num_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            if text_per_slide:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
                txBox.text_frame.text = text_per_slide

        path = Path(tempfile.mktemp(suffix=".pptx"))
        prs.save(str(path))
        return path

    def test_image_heavy_few_slides(self):
        """PPTX with < min_slides is not image-heavy regardless of text."""
        path = self._make_pptx(3, "")  # 3 slides, no text
        try:
            assert detect_image_heavy(path) is False
        finally:
            path.unlink()

    def test_image_heavy_no_text(self):
        """PPTX with many slides and no text is image-heavy."""
        path = self._make_pptx(10, "")
        try:
            assert detect_image_heavy(path) is True
        finally:
            path.unlink()

    def test_image_heavy_little_text(self):
        """PPTX with many slides and very little text is image-heavy."""
        path = self._make_pptx(10, "Hi")  # 2 chars * 10 = 20 chars total
        try:
            assert detect_image_heavy(path) is True
        finally:
            path.unlink()

    def test_not_image_heavy_with_text(self):
        """PPTX with substantial text per slide is not image-heavy."""
        long_text = "This is a detailed slide with lots of textual content. " * 5
        path = self._make_pptx(10, long_text)
        try:
            assert detect_image_heavy(path) is False
        finally:
            path.unlink()

    def test_custom_threshold(self):
        """Custom text_threshold changes detection sensitivity."""
        # "Short text" = 10 chars * 10 slides = 100 chars, 10 chars/slide
        path = self._make_pptx(10, "Short text")
        try:
            # With high threshold, should be image-heavy (100 < 5000)
            assert detect_image_heavy(path, text_threshold=5000) is True
            # With very low threshold AND raising chars/slide min
            # 100 chars > 50 threshold, but 10 chars/slide < 50 so still image-heavy
            # Need enough text to pass both checks
        finally:
            path.unlink()

        # Make a deck with enough text to pass both checks
        medium_text = "A" * 100  # 100 chars/slide
        path2 = self._make_pptx(10, medium_text)
        try:
            # 1000 chars total, 100 chars/slide — not image-heavy
            assert detect_image_heavy(path2, text_threshold=500) is False
        finally:
            path2.unlink()


# ── RenderedSlide dataclass tests ────────────────────────────────


class TestRenderedSlide:
    """Test RenderedSlide dataclass."""

    def test_fields(self):
        rs = RenderedSlide(
            slide_number=1,
            image_path=Path("/tmp/slide_001.png"),
            width=1920,
            height=1080,
        )
        assert rs.slide_number == 1
        assert rs.width == 1920


# ── MAX_SLIDES constant ─────────────────────────────────────────


class TestConstants:
    """Test module constants."""

    def test_max_slides_cap(self):
        """Max slides should be 50."""
        assert MAX_SLIDES == 50


# ── can_render tests ─────────────────────────────────────────────


class TestCanRender:
    """Test renderer availability detection."""

    def test_can_render_with_comtypes(self):
        """can_render returns True when comtypes is available."""
        with patch.dict("sys.modules", {"comtypes": MagicMock(), "comtypes.client": MagicMock()}):
            assert can_render() is True

    def test_can_render_no_backends(self):
        """can_render returns False when no backends available."""
        with patch("src.slides.renderer._find_libreoffice", return_value=None):
            with patch.dict("sys.modules", {"comtypes": None, "comtypes.client": None}):
                # Force ImportError for comtypes
                import builtins

                real_import = builtins.__import__

                def mock_import(name, *args, **kwargs):
                    if name in ("comtypes", "comtypes.client"):
                        raise ImportError("No comtypes")
                    return real_import(name, *args, **kwargs)

                with patch("builtins.__import__", side_effect=mock_import):
                    assert can_render() is False


# ── Tier router integration tests ────────────────────────────────


class TestTierRouterPptx:
    """Test PPTX tier routing logic."""

    def _make_pptx(self, num_slides: int, text_per_slide: str = "") -> Path:
        """Create a minimal PPTX file for testing."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        for i in range(num_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            if text_per_slide:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
                txBox.text_frame.text = text_per_slide

        path = Path(tempfile.mktemp(suffix=".pptx"))
        prs.save(str(path))
        return path

    def test_text_heavy_pptx_stays_tier2(self):
        """PPTX with lots of text stays at Tier 2."""
        from src.tier_router import route_tier, Tier
        from src.inventory import SourceFile, FileType

        long_text = "Detailed content with many words per slide. " * 10
        path = self._make_pptx(10, long_text)
        try:
            sf = SourceFile(path=path, type=FileType.SLIDES, size_bytes=path.stat().st_size, name=path.stem)
            decision = route_tier(sf)
            assert decision.tier == Tier.TEXT_AI
        finally:
            path.unlink()

    def test_image_heavy_pptx_routes_tier3_if_renderer(self):
        """Image-heavy PPTX routes to Tier 3 when renderer is available."""
        from src.tier_router import route_tier, Tier
        from src.inventory import SourceFile, FileType

        path = self._make_pptx(10, "")  # no text
        try:
            sf = SourceFile(path=path, type=FileType.SLIDES, size_bytes=path.stat().st_size, name=path.stem)
            with patch("src.slides.renderer.can_render", return_value=True):
                decision = route_tier(sf)
                assert decision.tier == Tier.MULTIMODAL
        finally:
            path.unlink()

    def test_image_heavy_pptx_falls_back_tier2_no_renderer(self):
        """Image-heavy PPTX falls back to Tier 2 when no renderer available."""
        from src.tier_router import route_tier, Tier
        from src.inventory import SourceFile, FileType

        path = self._make_pptx(10, "")
        try:
            sf = SourceFile(path=path, type=FileType.SLIDES, size_bytes=path.stat().st_size, name=path.stem)
            with patch("src.slides.renderer.can_render", return_value=False):
                decision = route_tier(sf)
                assert decision.tier == Tier.TEXT_AI
        finally:
            path.unlink()
