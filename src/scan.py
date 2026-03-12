"""
Tier 1 local metadata scanner — no API calls, fully local.

Scans files and extracts structured metadata using local libraries only.
Designed for corp-by-os `corp overnight --scope full-reshape` pre-scan.

Usage:
    from src.scan import scan_path, ScanResult

    results = scan_path(Path("C:/Users/docs"), recursive=True)
    print(f"Scanned {len(results)} files")
"""

import hashlib
import json
import logging
import subprocess
from dataclasses import dataclass, field, asdict
from datetime import datetime
from pathlib import Path

log = logging.getLogger(__name__)

# Max file size for SHA-256 hashing (500MB)
MAX_HASH_SIZE = 500 * 1024 * 1024

# Extensions we know how to extract rich metadata from
RICH_EXTENSIONS = {".pptx", ".pdf", ".docx", ".xlsx", ".csv", ".txt", ".md", ".rst", ".log"}
VIDEO_EXTENSIONS = {".mp4", ".mkv", ".webm", ".avi", ".mov", ".wmv"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".ogg", ".flac"}


@dataclass
class FileScanResult:
    """Scan result for a single file."""
    path: str
    filename: str
    extension: str
    size_bytes: int
    file_hash: str | None = None
    tier: int = 1
    metadata: dict = field(default_factory=dict)
    error: str | None = None


def _sha256_file(path: Path) -> str | None:
    """Compute SHA-256 hash of a file, reading in 64KB chunks. Skip files >500MB."""
    if path.stat().st_size > MAX_HASH_SIZE:
        return None
    h = hashlib.sha256()
    try:
        with open(path, "rb") as f:
            while True:
                chunk = f.read(65536)
                if not chunk:
                    break
                h.update(chunk)
        return f"sha256:{h.hexdigest()}"
    except (OSError, PermissionError) as exc:
        log.warning("Cannot hash %s: %s", path.name, exc)
        return None


def _scan_pptx(path: Path) -> dict:
    """Extract metadata from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slide_count = len(prs.slides)

    all_text_parts: list[str] = []
    notes_parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        all_text_parts.append(text)
        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes_parts.append(f"[Slide {i}] {notes_text}")

    full_text = "\n".join(all_text_parts)
    title = all_text_parts[0] if all_text_parts else path.stem

    meta: dict = {
        "title": title,
        "slide_count": slide_count,
        "text_chars": len(full_text),
        "text_preview": full_text[:2000],
    }
    if notes_parts:
        meta["speaker_notes_preview"] = "\n".join(notes_parts)[:1000]

    return meta


def _scan_pdf(path: Path) -> dict:
    """Extract metadata from PDF using pdfplumber."""
    import pdfplumber

    with pdfplumber.open(path) as pdf:
        page_count = len(pdf.pages)
        pages_text: list[str] = []
        # Extract text from first 3 pages for preview
        for page in pdf.pages[:3]:
            try:
                text = page.extract_text()
            except Exception:
                text = None
            if text:
                pages_text.append(text)

        # Try to get title from metadata
        info = pdf.metadata or {}
        title = info.get("Title") or None

    full_preview = "\n\n".join(pages_text)
    return {
        "title": title,
        "page_count": page_count,
        "text_chars": len(full_preview),
        "text_preview": full_preview[:2000],
    }


def _scan_docx(path: Path) -> dict:
    """Extract metadata from DOCX using python-docx."""
    import docx

    doc = docx.Document(str(path))

    title = None
    headings: list[str] = []
    text_parts: list[str] = []

    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        text_parts.append(text)
        if p.style and p.style.name:
            if p.style.name.startswith("Heading"):
                headings.append(text)
            if p.style.name == "Title" and title is None:
                title = text

    if not title and headings:
        title = headings[0]

    full_text = "\n".join(text_parts)
    return {
        "title": title,
        "headings": headings[:20],
        "text_chars": len(full_text),
        "text_preview": full_text[:2000],
    }


def _scan_xlsx(path: Path) -> dict:
    """Extract metadata from XLSX using openpyxl."""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheets: list[dict] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_data: list[list[str]] = []
        for row in ws.iter_rows(values_only=True, max_row=4):
            if row is None:
                continue
            cells = [str(c) if c is not None else "" for c in row]
            if any(c for c in cells):
                rows_data.append(cells)

        sheet_info: dict = {"name": sheet_name}
        if rows_data:
            sheet_info["headers"] = rows_data[0]
            if len(rows_data) > 1:
                sheet_info["sample_rows"] = rows_data[1:4]

        sheets.append(sheet_info)

    wb.close()
    return {
        "sheet_names": wb.sheetnames,
        "sheet_count": len(wb.sheetnames),
        "sheets": sheets,
    }


def _scan_csv(path: Path) -> dict:
    """Extract metadata from CSV."""
    import csv

    rows: list[list[str]] = []
    row_count = 0
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                row_count += 1
                if len(rows) < 4:
                    rows.append(row)
    except Exception as exc:
        return {"error": str(exc)}

    meta: dict = {"row_count": row_count}
    if rows:
        meta["headers"] = rows[0]
        if len(rows) > 1:
            meta["sample_rows"] = rows[1:4]
    return meta


def _scan_text(path: Path) -> dict:
    """Extract metadata from plain text files."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return {"error": str(exc)}
    return {
        "text_chars": len(text),
        "text_preview": text[:2000],
        "line_count": text.count("\n") + 1,
    }


def _scan_video(path: Path) -> dict:
    """Extract metadata from video using ffprobe (if available)."""
    meta: dict = {}
    try:
        result = subprocess.run(
            [
                "ffprobe", "-v", "quiet",
                "-print_format", "json",
                "-show_format", "-show_streams",
                str(path),
            ],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0:
            probe = json.loads(result.stdout)
            fmt = probe.get("format", {})
            duration = fmt.get("duration")
            if duration:
                meta["duration_seconds"] = round(float(duration))

            # Find video stream
            for stream in probe.get("streams", []):
                if stream.get("codec_type") == "video":
                    meta["codec"] = stream.get("codec_name")
                    w = stream.get("width")
                    h = stream.get("height")
                    if w and h:
                        meta["resolution"] = f"{w}x{h}"
                    break
    except FileNotFoundError:
        log.debug("ffprobe not available, returning size-only metadata for %s", path.name)
    except (subprocess.TimeoutExpired, json.JSONDecodeError, ValueError) as exc:
        log.warning("ffprobe failed for %s: %s", path.name, exc)
    return meta


def _determine_tier(ext: str, metadata: dict) -> int:
    """Determine extraction tier based on file type."""
    if ext in VIDEO_EXTENSIONS or ext in AUDIO_EXTENSIONS:
        return 3
    if ext in RICH_EXTENSIONS:
        return 1 if metadata.get("text_chars", 0) > 0 else 2
    return 1


def scan_file(path: Path, base_path: Path | None = None) -> FileScanResult:
    """Scan a single file and extract local metadata.

    Args:
        path: Absolute path to file
        base_path: If set, store relative path from this base

    Returns:
        FileScanResult with metadata
    """
    ext = path.suffix.lower()
    rel_path = str(path.relative_to(base_path)).replace("\\", "/") if base_path else str(path).replace("\\", "/")

    result = FileScanResult(
        path=rel_path,
        filename=path.name,
        extension=ext,
        size_bytes=path.stat().st_size,
    )

    # Hash
    result.file_hash = _sha256_file(path)

    # Extract rich metadata based on extension
    try:
        if ext == ".pptx":
            result.metadata = _scan_pptx(path)
        elif ext == ".pdf":
            result.metadata = _scan_pdf(path)
        elif ext == ".docx":
            result.metadata = _scan_docx(path)
        elif ext == ".xlsx" or ext == ".xlsm":
            result.metadata = _scan_xlsx(path)
        elif ext == ".csv":
            result.metadata = _scan_csv(path)
        elif ext in (".txt", ".md", ".markdown", ".rst", ".log"):
            result.metadata = _scan_text(path)
        elif ext in VIDEO_EXTENSIONS:
            result.metadata = _scan_video(path)
        elif ext in AUDIO_EXTENSIONS:
            result.metadata = _scan_video(path)  # ffprobe works for audio too
    except Exception as exc:
        log.warning("Metadata extraction failed for %s: %s", path.name, exc)
        result.error = str(exc)

    result.tier = _determine_tier(ext, result.metadata)
    return result


def scan_path(
    input_path: Path,
    recursive: bool = True,
    exclude: tuple[str, ...] = (),
) -> list[FileScanResult]:
    """Scan a file or directory and return metadata for all files.

    Args:
        input_path: Path to file or directory
        recursive: Whether to scan subdirectories
        exclude: Folder names to skip

    Returns:
        List of FileScanResult objects
    """
    results: list[FileScanResult] = []
    exclude_set = set(exclude)

    if input_path.is_file():
        results.append(scan_file(input_path))
        return results

    if not input_path.is_dir():
        raise ValueError(f"Path does not exist: {input_path}")

    base_path = input_path
    pattern = "**/*" if recursive else "*"
    all_files = sorted(input_path.glob(pattern))

    count = 0
    for p in all_files:
        if not p.is_file():
            continue
        # Skip hidden files and folders
        if any(part.startswith(".") for part in p.relative_to(base_path).parts):
            continue
        # Skip excluded folders
        if any(part in exclude_set for part in p.relative_to(base_path).parts):
            continue

        count += 1
        if count % 100 == 0:
            log.info("Scanned %d files...", count)

        results.append(scan_file(p, base_path))

    log.info("Scan complete: %d files", count)
    return results


def results_to_json(results: list[FileScanResult]) -> dict:
    """Convert scan results to the output JSON structure."""
    return {
        "scan_date": datetime.now().isoformat(),
        "total_files": len(results),
        "results": [asdict(r) for r in results],
    }
