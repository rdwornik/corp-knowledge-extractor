"""
Local text extraction from documents — Tier 1 (FREE, no API calls).

Extracts plain text from PDF, DOCX, PPTX, XLSX, CSV, and text files
using local libraries. Returns structured result with quality metrics
so the tier router can decide whether AI extraction is needed.

Usage:
    from src.text_extract import extract_text, TextExtractionResult

    result = extract_text(Path("report.pdf"))
    if result.extraction_quality == "good":
        # Tier 1 or 2 is sufficient
    else:
        # Fall back to Tier 3 multimodal
"""

import logging
from dataclasses import dataclass, field
from pathlib import Path

log = logging.getLogger(__name__)


@dataclass
class TextExtractionResult:
    """Result of local text extraction."""
    text: str
    char_count: int = 0
    page_count: int = 0
    slide_count: int = 0
    has_images: bool = False
    extraction_quality: str = "none"  # "good", "partial", "none"
    extractor: str = "unknown"
    error: str | None = None


def extract_source_date(file_path: Path) -> str | None:
    """Extract document date from file metadata.

    Returns YYYY-MM or YYYY-MM-DD, or None if unavailable.
    Uses file metadata only — never content or LLM.
    """
    ext = file_path.suffix.lower()
    try:
        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                info = pdf.metadata or {}
                date_str = info.get("CreationDate") or info.get("ModDate")
                if date_str and isinstance(date_str, str) and date_str.startswith("D:"):
                    return f"{date_str[2:6]}-{date_str[6:8]}"
                return None
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            mod = prs.core_properties.modified
            if mod:
                return mod.strftime("%Y-%m")
            return None
        elif ext == ".docx":
            import docx
            doc = docx.Document(str(file_path))
            mod = doc.core_properties.modified
            if mod:
                return mod.strftime("%Y-%m")
            return None
        elif ext in (".xlsx", ".xlsm"):
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True)
            mod = wb.properties.modified
            wb.close()
            if mod:
                return mod.strftime("%Y-%m")
            return None
    except Exception:
        return None
    return None


def extract_text(path: Path) -> TextExtractionResult:
    """Extract text from a file using the appropriate local extractor.

    Args:
        path: Path to the file to extract text from

    Returns:
        TextExtractionResult with extracted text and quality metrics
    """
    ext = path.suffix.lower()
    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
        ".csv": _extract_csv,
        ".txt": _extract_plaintext,
        ".md": _extract_plaintext,
        ".markdown": _extract_plaintext,
        ".rst": _extract_plaintext,
        ".log": _extract_plaintext,
    }

    extractor = extractors.get(ext)
    if not extractor:
        return TextExtractionResult(
            text="",
            extraction_quality="none",
            extractor="unsupported",
            error=f"No local extractor for {ext}",
        )

    try:
        return extractor(path)
    except Exception as exc:
        log.warning("Local text extraction failed for %s: %s", path.name, exc)
        return TextExtractionResult(
            text="",
            extraction_quality="none",
            extractor=ext.lstrip("."),
            error=str(exc),
        )


def _assess_quality(text: str, page_count: int = 0) -> str:
    """Assess extraction quality based on text density."""
    char_count = len(text)
    if char_count == 0:
        return "none"
    # Very short text relative to page count suggests scanned/image-heavy doc
    if page_count > 0 and char_count / page_count < 100:
        return "partial"
    if char_count < 200:
        return "partial"
    return "good"


def _extract_pdf(path: Path) -> TextExtractionResult:
    """Extract text from PDF using pdfplumber. Includes [PAGE N] markers."""
    import pdfplumber

    pages_text = []
    has_images = False
    had_none_pages = False

    with pdfplumber.open(path) as pdf:
        page_count = len(pdf.pages)
        for i, page in enumerate(pdf.pages, 1):
            try:
                text = page.extract_text()
            except Exception:
                text = None
            if text is None:
                had_none_pages = True
                text = ""
            pages_text.append(f"[PAGE {i}]\n{text}")
            if page.images:
                has_images = True

    full_text = "\n\n".join(pages_text)
    quality = _assess_quality(full_text, page_count)
    # Downgrade to partial if some pages returned None
    if had_none_pages and quality == "good":
        quality = "partial"
    return TextExtractionResult(
        text=full_text,
        char_count=len(full_text),
        page_count=page_count,
        has_images=has_images,
        extraction_quality=quality,
        extractor="pdfplumber",
    )


def _extract_docx(path: Path) -> TextExtractionResult:
    """Extract text from DOCX using python-docx. Includes [SECTION N] markers."""
    import docx

    doc = docx.Document(path)
    has_images = any(
        r.element.tag.endswith("}drawing") or r.element.tag.endswith("}pict")
        for p in doc.paragraphs
        for r in p.runs
        for child in r.element
    )

    # Group paragraphs into sections by headings
    sections: list[str] = []
    current: list[str] = []
    section_num = 1
    for p in doc.paragraphs:
        if not p.text.strip():
            continue
        if p.style and p.style.name and p.style.name.startswith("Heading"):
            if current:
                sections.append(f"[SECTION {section_num}]\n" + "\n\n".join(current))
                section_num += 1
                current = []
        current.append(p.text)
    if current:
        sections.append(f"[SECTION {section_num}]\n" + "\n\n".join(current))

    full_text = "\n\n".join(sections)
    return TextExtractionResult(
        text=full_text,
        char_count=len(full_text),
        page_count=0,  # DOCX doesn't have reliable page count
        has_images=has_images,
        extraction_quality=_assess_quality(full_text),
        extractor="python-docx",
    )


def _extract_pptx(path: Path) -> TextExtractionResult:
    """Extract text from PPTX using python-pptx. Includes [SLIDE N] markers."""
    from pptx import Presentation

    prs = Presentation(path)
    slides_text = []
    has_images = False

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                has_images = True
        slide_body = "\n".join(texts) if texts else ""
        slides_text.append(f"[SLIDE {i}]\n{slide_body}")

    full_text = "\n\n---\n\n".join(slides_text)
    return TextExtractionResult(
        text=full_text,
        char_count=len(full_text),
        slide_count=len(prs.slides),
        has_images=has_images,
        extraction_quality=_assess_quality(full_text, len(prs.slides)),
        extractor="python-pptx",
    )


def _extract_xlsx(path: Path) -> TextExtractionResult:
    """Extract text from XLSX using openpyxl."""
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            if row is None:
                continue
            cells = [str(c) if c is not None else "" for c in row]
            joined = " | ".join(c for c in cells if c)
            if joined:
                rows.append(joined)
        if rows:
            sheets_text.append(f"## {sheet_name}\n" + "\n".join(rows))

    wb.close()
    full_text = "\n\n".join(sheets_text)
    return TextExtractionResult(
        text=full_text,
        char_count=len(full_text),
        extraction_quality=_assess_quality(full_text),
        extractor="openpyxl",
    )


def _extract_csv(path: Path) -> TextExtractionResult:
    """Extract text from CSV."""
    import csv

    rows = []
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(row))

    full_text = "\n".join(rows)
    return TextExtractionResult(
        text=full_text,
        char_count=len(full_text),
        extraction_quality=_assess_quality(full_text),
        extractor="csv",
    )


def _extract_plaintext(path: Path) -> TextExtractionResult:
    """Extract text from plain text files."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return TextExtractionResult(
        text=text,
        char_count=len(text),
        extraction_quality=_assess_quality(text),
        extractor="plaintext",
    )
