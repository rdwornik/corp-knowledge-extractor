"""Render PPTX slides as PNG images.

Primary: PowerPoint COM automation (Windows, best quality)
Fallback: LibreOffice headless -> PDF -> PyMuPDF -> PNG (cross-platform)

If neither renderer is available, raises RuntimeError so the caller
can fall back to Tier 2 text-only extraction.
"""

import logging
import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

# Max slides to render (cost control — each slide is ~0.001 in vision tokens)
MAX_SLIDES = 50


@dataclass
class RenderedSlide:
    """A single rendered slide image."""

    slide_number: int
    image_path: Path
    width: int
    height: int


def render_slides(pptx_path: Path, output_dir: Path) -> list[RenderedSlide]:
    """Render all slides from PPTX as PNG images.

    Tries PowerPoint COM first (Windows), falls back to LibreOffice.
    Caps at MAX_SLIDES slides.

    Returns:
        List of RenderedSlide with paths to PNG files.

    Raises:
        RuntimeError: If no renderer is available.
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    # Try PowerPoint COM (Windows only, best quality)
    try:
        slides = _render_via_powerpoint(pptx_path, output_dir)
        if slides:
            return slides[:MAX_SLIDES]
    except Exception as e:
        logger.info("PowerPoint COM not available: %s, trying LibreOffice", e)

    # Try LibreOffice headless
    try:
        slides = _render_via_libreoffice(pptx_path, output_dir)
        if slides:
            return slides[:MAX_SLIDES]
    except Exception as e:
        logger.info("LibreOffice not available: %s", e)

    raise RuntimeError("Cannot render PPTX slides. Install PowerPoint (Windows) or LibreOffice.")


def _render_via_powerpoint(pptx_path: Path, output_dir: Path) -> list[RenderedSlide]:
    """Render via PowerPoint COM automation. Windows only."""
    import comtypes.client

    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(
        str(pptx_path.resolve()),
        ReadOnly=True,
        Untitled=False,
        WithWindow=False,
    )

    slides = []
    try:
        for i, slide in enumerate(presentation.Slides, 1):
            png_path = output_dir / f"slide_{i:03d}.png"
            slide.Export(str(png_path.resolve()), "PNG", 1920, 1080)
            slides.append(
                RenderedSlide(
                    slide_number=i,
                    image_path=png_path,
                    width=1920,
                    height=1080,
                )
            )
            if i >= MAX_SLIDES:
                logger.warning("Capped at %d slides", MAX_SLIDES)
                break
    finally:
        presentation.Close()

    logger.info("Rendered %d slides via PowerPoint COM", len(slides))
    return slides


def _render_via_libreoffice(pptx_path: Path, output_dir: Path) -> list[RenderedSlide]:
    """Render via LibreOffice headless -> PDF -> PyMuPDF -> PNG."""
    lo_path = _find_libreoffice()
    if not lo_path:
        raise FileNotFoundError("LibreOffice not found")

    # Convert PPTX -> PDF
    pdf_dir = output_dir / "_pdf_temp"
    pdf_dir.mkdir(exist_ok=True)

    subprocess.run(
        [lo_path, "--headless", "--convert-to", "pdf", "--outdir", str(pdf_dir), str(pptx_path)],
        check=True,
        timeout=120,
        capture_output=True,
    )

    pdf_path = pdf_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF not generated: {pdf_path}")

    # PDF -> PNG per page via PyMuPDF
    import fitz

    doc = fitz.open(str(pdf_path))
    slides = []

    for i, page in enumerate(doc, 1):
        mat = fitz.Matrix(2, 2)  # 2x zoom for quality (144 DPI)
        pix = page.get_pixmap(matrix=mat)
        png_path = output_dir / f"slide_{i:03d}.png"
        pix.save(str(png_path))
        slides.append(
            RenderedSlide(
                slide_number=i,
                image_path=png_path,
                width=pix.width,
                height=pix.height,
            )
        )
        if i >= MAX_SLIDES:
            logger.warning("Capped at %d slides", MAX_SLIDES)
            break

    doc.close()

    # Cleanup temp PDF
    pdf_path.unlink(missing_ok=True)
    try:
        pdf_dir.rmdir()
    except OSError:
        pass

    logger.info("Rendered %d slides via LibreOffice + PyMuPDF", len(slides))
    return slides


def _find_libreoffice() -> str | None:
    """Find LibreOffice executable."""
    # PATH lookup
    for name in ("libreoffice", "soffice"):
        found = shutil.which(name)
        if found:
            return found

    # Common Windows paths
    for candidate in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]:
        if Path(candidate).exists():
            return candidate

    return None


def detect_image_heavy(
    pptx_path: Path,
    text_threshold: int = 500,
    min_slides: int = 5,
) -> bool:
    """Heuristic: is this PPTX image-heavy and needs multimodal?

    Returns True if total extracted text < text_threshold AND slide count >= min_slides.
    This catches diagram-heavy decks where python-pptx gets almost no text.
    """
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slide_count = len(prs.slides)

    if slide_count < min_slides:
        return False

    total_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                total_text += shape.text_frame.text + " "

    text_len = len(total_text.strip())
    text_per_slide = text_len / max(slide_count, 1)
    is_heavy = text_len < text_threshold or text_per_slide < 50

    if is_heavy:
        logger.info(
            "PPTX image-heavy: %d chars across %d slides (%.0f chars/slide)",
            text_len,
            slide_count,
            text_per_slide,
        )

    return is_heavy


def can_render() -> bool:
    """Check if any slide renderer is available (without actually rendering)."""
    # Check comtypes (PowerPoint COM)
    try:
        import comtypes.client  # noqa: F401

        return True
    except ImportError:
        pass

    # Check LibreOffice + PyMuPDF
    if _find_libreoffice():
        try:
            import fitz  # noqa: F401

            return True
        except ImportError:
            pass

    return False
